# Mohammed Alshehhi
# 31/10/2023
# AI Quiz Generator

import os
from pptx import Presentation
import openai
import tkinter as tk
from tkinter import filedialog, Text, messagebox

# Set your OpenAI API key here
openai.api_key = "API_Key_Here"


# Function to load a PowerPoint file
def load_pptx():
    global pptx_path
    pptx_path = filedialog.askopenfilename(filetypes=[("PowerPoint files", "*.pptx")])
    if pptx_path:
        file_label.config(text=os.path.basename(pptx_path))


# Function to generate quiz questions
def generate_questions():
    # Define the types of questions and the number to generate
    question_types = {
        "mcq": ("multiple choice questions", int(mcq_var.get())),
        "short": ("short answers questions", int(short_var.get())),
        "tf": ("true/false questions", int(tf_var.get())),
        "fitb": ("fill-in-the-blank questions", int(fitb_var.get())),
        "scenario": ("scenario based questions", int(scenario_var.get())),
        "other": (other_var.get(), 1 if other_var.get() else 0)
    }

    # Initialize the messages list for the chat completion API call
    messages = [
        {"role": "system",
         "content": "you will receive slides from power point, then you will generate a quiz based on the requirements"},
    ]

    # Load the PowerPoint presentation
    pr = Presentation(pptx_path)
    # Iterate over each slide to extract text
    for i, slide in enumerate(pr.slides):
        slide_text = ''
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text += shape.text
        sn = i + 1
        messages.append({"role": "user", "content": "Slide " + str(sn) + ":\n\n" + slide_text})

    # Clear the output text box before inserting new content
    output_text.delete(1.0, tk.END)
    # Mark the end of slides
    messages.append({"role": "user", "content": "End of Slides"})

    # Generate questions for each specified type
    for question_type, (question_text, num_questions) in question_types.items():
        if num_questions > 0:
            output_text.insert(tk.END, "-------------  " + question_text + "  -------------\n\n")
            if question_type == "mcq":
                mcq_label.config(bg="orange")
            if question_type == "short":
                short_label.config(bg="orange")
            if question_type == "tf":
                tf_label.config(bg="orange")
            if question_type == "fitb":
                fitb_label.config(bg="orange")
            if question_type == "scenario":
                scenario_label.config(bg="orange")
            if question_type == "other":
                other_label.config(bg="orange")
            root.update()
            root.update()
            root.title("Generating " + question_text)

            # Prepare the content request for the AI
            request_content = f"and now, from the slides, generate {num_questions} {question_text}"
            request_content += ". Make sure it follows all Cognitive Complexity levels"
            if question_text == "multiple choice questions":
                request_content += ". \nMake sure if follows the following : Use clear, straightforward language in " \
                                   "both stem/options" \
                                   "\nBe careful of grammar that might give away the key" \
                                   "\nAvoid negatives in the stem; 'none / all of the above' in the options" \
                                   "\nHighlight NOT, EXCEPT, TRUE, FALSE, MOST, BEST if you use them in the stem" \
                                   "\nAvoid repetitive words in the options" \
                                   "\nTry to keep the language in the options parallel" \
                                   "\nMake the distractors plausible; avoid humorous answers" \
                                   "\nMake sure your test item is valid: does it measure what it is supposed to measure?"
            messages.append({"role": "user", "content": request_content})

            # Call the OpenAI API to generate questions
            response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo-16k",
                messages=messages
            )

            # Insert the AI-generated questions into the output text box
            output_text.insert(tk.END, response["choices"][0]["message"]["content"])
            messages.append({"role": "assistant", "content": response["choices"][0]["message"]["content"]})
            output_text.insert(tk.END, "\n\n")
            if question_type == "mcq":
                mcq_label.config(bg="#25BE2F")
            if question_type == "short":
                short_label.config(bg="#25BE2F")
            if question_type == "tf":
                tf_label.config(bg="#25BE2F")
            if question_type == "fitb":
                fitb_label.config(bg="#25BE2F")
            if question_type == "scenario":
                scenario_label.config(bg="#25BE2F")
            if question_type == "other":
                other_label.config(bg="#25BE2F")
            root.update()
            root.update()
            root.update()
    root.title("Done!")


def copy_output():
    root.clipboard_clear()
    root.clipboard_append(output_text.get(1.0, tk.END))


root = tk.Tk()
root.title("Quiz Generator")

frame = tk.Frame(root, padx=20, pady=20)
frame.pack()

load_btn = tk.Button(frame, text="Load PowerPoint", command=load_pptx)
load_btn.grid(row=0, column=0, pady=10)

file_label = tk.Label(frame, text="No file selected")
file_label.grid(row=0, column=1, pady=10)

mcq_label = tk.Label(frame, text="Multiple Choice Questions:")
mcq_label.grid(row=1, column=0)
mcq_var = tk.Entry(frame, width=5)
mcq_var.insert(0, "0")
mcq_var.grid(row=1, column=1)

short_label = tk.Label(frame, text="Short Answer Questions:")
short_label.grid(row=2, column=0)
short_var = tk.Entry(frame, width=5)
short_var.insert(0, "0")
short_var.grid(row=2, column=1)

tf_label = tk.Label(frame, text="True/False Questions:")
tf_label.grid(row=3, column=0)
tf_var = tk.Entry(frame, width=5)
tf_var.insert(0, "0")
tf_var.grid(row=3, column=1)

fitb_label = tk.Label(frame, text="Fill-in-the-blank Questions:")
fitb_label.grid(row=4, column=0)
fitb_var = tk.Entry(frame, width=5)
fitb_var.insert(0, "0")
fitb_var.grid(row=4, column=1)

scenario_label = tk.Label(frame, text="Scenario-based Questions:")
scenario_label.grid(row=5, column=0)
scenario_var = tk.Entry(frame, width=5)
scenario_var.insert(0, "0")
scenario_var.grid(row=5, column=1)

other_label = tk.Label(frame, text="Other (custom):")
other_label.grid(row=6, column=0)
other_var = tk.Entry(frame, width=15)
other_var.grid(row=6, column=1)

generate_btn = tk.Button(frame, text="Generate Questions", command=generate_questions)
generate_btn.grid(row=7, column=0, columnspan=2, pady=10)

output_label = tk.Label(frame, text="Output:")
output_label.grid(row=8, column=0, sticky="nw")

output_text = Text(frame, width=60, height=20, wrap="word")
output_text.grid(row=9, column=0, columnspan=2)

scrollbar = tk.Scrollbar(frame, orient="vertical", command=output_text.yview)
scrollbar.grid(row=9, column=2, sticky="ns")
output_text.config(yscrollcommand=scrollbar.set)

copy_btn = tk.Button(frame, text="Copy to Clipboard", command=copy_output)
copy_btn.grid(row=10, column=0, columnspan=2, pady=10)

root.mainloop()
